"""FinEdit 営業資料 .pptx 生成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ---- Brand colors ----
NAVY       = RGBColor(0x0a, 0x19, 0x2f)
NAVY_LIGHT = RGBColor(0x17, 0x2a, 0x45)
GOLD       = RGBColor(0xd4, 0xaf, 0x37)
WHITE      = RGBColor(0xff, 0xff, 0xff)
GRAY_BG    = RGBColor(0xf5, 0xf6, 0xf8)
TEXT_DARK  = RGBColor(0x1a, 0x1a, 0x2e)
TEXT_MID   = RGBColor(0x55, 0x60, 0x70)

# ---- Slide size: widescreen 16:9 ----
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Noto Sans JP"
    return txBox


def add_label(slide, text, x, y):
    """Small gold uppercase label above section heading."""
    add_text(slide, text, x, y, Inches(6), Inches(0.35),
             size=10, bold=True, color=GOLD, align=PP_ALIGN.LEFT)


def add_heading(slide, text, x, y, w=Inches(10), dark_bg=True):
    c = WHITE if dark_bg else TEXT_DARK
    add_text(slide, text, x, y, w, Inches(0.9),
             size=28, bold=True, color=c, align=PP_ALIGN.LEFT)


def add_divider(slide, x, y, w=Inches(1.2)):
    add_rect(slide, x, y, w, Pt(3), fill=GOLD)


# ============================================================
# Slide 01 – Cover
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=NAVY)
# Gold accent bar left
add_rect(sl, 0, 0, Inches(0.12), H, fill=GOLD)
# Gold accent bar bottom
add_rect(sl, 0, H - Inches(0.08), W, Inches(0.08), fill=GOLD)

# Eyebrow
add_text(sl, "名古屋の経営者・起業家向け AI活用支援",
         Inches(1.2), Inches(1.4), Inches(10), Inches(0.45),
         size=13, color=GOLD)

# Main title
add_text(sl, "AIを、使える知識から",
         Inches(1.2), Inches(2.0), Inches(11), Inches(0.75),
         size=36, bold=True, color=WHITE)
add_text(sl, "「仕事が回る仕組み」へ。",
         Inches(1.2), Inches(2.75), Inches(11), Inches(0.75),
         size=36, bold=True, color=GOLD)

# Sub copy
add_text(sl,
         "ChatGPT・Gemini・Claudeなどを実務に落とし込み、\n"
         "ホームページ・業務フロー・SNS発信まで一緒に整える伴走支援です。",
         Inches(1.2), Inches(3.7), Inches(9), Inches(1.1),
         size=14, color=RGBColor(0xcc, 0xd6, 0xf1))

# Tag pill (simulated)
add_rect(sl, Inches(1.2), Inches(5.0), Inches(4.2), Inches(0.42),
         fill=NAVY_LIGHT, line=GOLD, line_w=Pt(1))
add_text(sl, "AI活用インストラクター / 実装パートナー",
         Inches(1.3), Inches(5.02), Inches(4.0), Inches(0.38),
         size=11, color=GOLD, align=PP_ALIGN.CENTER)

# Brand / name
add_text(sl, "FinEdit  |  勝目 麻希",
         Inches(1.2), Inches(6.4), Inches(5), Inches(0.45),
         size=12, color=RGBColor(0x88, 0x92, 0xb0))

# Slide number
add_text(sl, "01 / 09", W - Inches(1.5), H - Inches(0.45), Inches(1.4), Inches(0.35),
         size=9, color=RGBColor(0x44, 0x50, 0x60), align=PP_ALIGN.RIGHT)


# ============================================================
# Slide 02 – Problems
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GRAY_BG)
add_rect(sl, 0, 0, W, Inches(0.08), fill=GOLD)

add_label(sl, "PROBLEM", Inches(0.7), Inches(0.6))
add_heading(sl, "こんなお悩みはありませんか？",
            Inches(0.7), Inches(0.95), dark_bg=False)
add_divider(sl, Inches(0.7), Inches(1.75))

problems = [
    ("ChatGPTを使ってみたが、\n仕事の成果につながっていない",),
    ("どの業務にAIを使えばよいか、\n社内で判断できる人がいない",),
    ("情報漏洩やセキュリティが不安で、\n社内導入に踏み切れない",),
    ("AI研修を受けたが、\n日常業務に定着していない",),
    ("ホームページを作りたいが、\n何から始めればいいかわからない",),
    ("SNSや情報発信を始めたいが、\n継続できる仕組みがない",),
]

cols, rows = 3, 2
cw, ch = Inches(3.7), Inches(1.65)
gap_x, gap_y = Inches(0.25), Inches(0.22)
start_x, start_y = Inches(0.65), Inches(2.0)

for i, (txt,) in enumerate(problems):
    c, r = i % cols, i // cols
    x = start_x + c * (cw + gap_x)
    y = start_y + r * (ch + gap_y)
    add_rect(sl, x, y, cw, ch, fill=WHITE,
             line=RGBColor(0xe0, 0xe4, 0xea), line_w=Pt(1))
    add_rect(sl, x, y, Inches(0.07), ch, fill=GOLD)
    add_text(sl, txt, x + Inches(0.18), y + Inches(0.22),
             cw - Inches(0.28), ch - Inches(0.35),
             size=12, color=TEXT_DARK, wrap=True)

add_text(sl, "02 / 09", W - Inches(1.5), H - Inches(0.45), Inches(1.4), Inches(0.35),
         size=9, color=TEXT_MID, align=PP_ALIGN.RIGHT)


# ============================================================
# Slide 03 – Strengths
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=NAVY)
add_rect(sl, 0, 0, W, Inches(0.08), fill=GOLD)

add_label(sl, "OUR STRENGTHS", Inches(0.7), Inches(0.6))
add_heading(sl, "FinEditが選ばれる 3つの理由", Inches(0.7), Inches(0.95))
add_divider(sl, Inches(0.7), Inches(1.75))

strengths = [
    ("01", "整える力",
     "元銀行員としての視点で、業務・数字・リスクを整理。\n安全で現実的なAI活用の入口を設計します。"),
    ("02", "伝える力",
     "金融・ビジネス領域のプロライターとして、\n事業の魅力や提供価値をわかりやすい言葉に変換します。"),
    ("03", "形にする力",
     "講座で終わらせず、LP・業務フロー・SNS導線・\n実務プロンプトまで、使える形に落とし込みます。"),
]

cw, ch = Inches(3.7), Inches(3.5)
gap_x = Inches(0.4)
start_x, start_y = Inches(0.67), Inches(2.1)

for i, (num, title, body) in enumerate(strengths):
    x = start_x + i * (cw + gap_x)
    add_rect(sl, x, start_y, cw, ch, fill=NAVY_LIGHT,
             line=RGBColor(0x2a, 0x3f, 0x58), line_w=Pt(1))
    add_rect(sl, x, start_y, cw, Inches(0.07), fill=GOLD)

    # big number
    add_text(sl, num, x + Inches(0.2), start_y + Inches(0.2),
             Inches(1.5), Inches(0.9),
             size=44, bold=True, color=RGBColor(0x2a, 0x3f, 0x58))
    # title
    add_text(sl, title, x + Inches(0.2), start_y + Inches(1.1),
             cw - Inches(0.35), Inches(0.55),
             size=20, bold=True, color=GOLD)
    # body
    add_text(sl, body, x + Inches(0.2), start_y + Inches(1.75),
             cw - Inches(0.35), Inches(1.6),
             size=11.5, color=RGBColor(0xcc, 0xd6, 0xf1), wrap=True)

add_text(sl, "03 / 09", W - Inches(1.5), H - Inches(0.45), Inches(1.4), Inches(0.35),
         size=9, color=RGBColor(0x44, 0x50, 0x60), align=PP_ALIGN.RIGHT)


# ============================================================
# Slide 04 – Can-Do
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
add_rect(sl, 0, 0, W, Inches(0.08), fill=GOLD)

add_label(sl, "WHAT WE DO", Inches(0.7), Inches(0.6))
add_heading(sl, "一緒に整えられるもの", Inches(0.7), Inches(0.95), dark_bg=False)
add_divider(sl, Inches(0.7), Inches(1.75))

cando = [
    ("サービスLP・ホームページ", "構成、文章、デザイン、コード生成まで\n一気通貫で整えます。"),
    ("業務フロー・マニュアル",  "属人化した業務を可視化し、\n誰でも回せる形に整理します。"),
    ("SNS・発信設計",          "投稿ネタ、構成、文章、プロフィール、\n導線をAIで仕組み化します。"),
    ("AI活用の仕組み化",       "型・手順・テンプレートを整え、\n社内で再現できる仕組みにします。"),
    ("簡易業務アプリ",         "顧客管理、日報、教材など\n小さく使えるツールづくりを支援します。"),
    ("企画書・提案資料",       "事業の強みを整理し、\n伝わる資料や営業導線に落とし込みます。"),
]

cw, ch = Inches(3.7), Inches(1.65)
gap_x, gap_y = Inches(0.25), Inches(0.22)
start_x, start_y = Inches(0.65), Inches(2.0)

for i, (title, body) in enumerate(cando):
    c, r = i % 3, i // 3
    x = start_x + c * (cw + gap_x)
    y = start_y + r * (ch + gap_y)
    add_rect(sl, x, y, cw, ch, fill=GRAY_BG,
             line=RGBColor(0xe0, 0xe4, 0xea), line_w=Pt(1))
    add_rect(sl, x, y, Inches(0.06), ch, fill=NAVY)
    add_text(sl, title, x + Inches(0.16), y + Inches(0.15),
             cw - Inches(0.25), Inches(0.45),
             size=12, bold=True, color=NAVY)
    add_text(sl, body, x + Inches(0.16), y + Inches(0.65),
             cw - Inches(0.25), Inches(0.9),
             size=10.5, color=TEXT_MID, wrap=True)

add_text(sl, "04 / 09", W - Inches(1.5), H - Inches(0.45), Inches(1.4), Inches(0.35),
         size=9, color=TEXT_MID, align=PP_ALIGN.RIGHT)


# ============================================================
# Slide 05 – Services
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=NAVY)
add_rect(sl, 0, 0, W, Inches(0.08), fill=GOLD)

add_label(sl, "SERVICES", Inches(0.7), Inches(0.6))
add_heading(sl, "サービス内容", Inches(0.7), Inches(0.95))
add_divider(sl, Inches(0.7), Inches(1.75))

services = [
    ("01", "AIホームページ・LP制作支援",
     "事業内容を整理し、集客導線まで考えたホームページ・LPをAIと一緒に形にします。"),
    ("02", "業務フロー整理・AI活用設計",
     "繰り返し業務や属人化している作業を可視化し、AIで効率化できる形へ落とし込みます。"),
    ("03", "SNS・発信サポート",
     "投稿ネタ、構成、文章、プロフィール、サービス導線まで発信全体を整えます。"),
    ("04", "AI活用研修・ワークショップ",
     "AI未経験の方にもわかりやすく、当日から使えるプロンプトと運用ルールを設計します。"),
]

sw, sh = Inches(11.5), Inches(1.1)
start_x, start_y = Inches(0.7), Inches(2.1)
gap_y = Inches(0.18)

for i, (num, title, body) in enumerate(services):
    y = start_y + i * (sh + gap_y)
    add_rect(sl, start_x, y, sw, sh, fill=NAVY_LIGHT,
             line=RGBColor(0x2a, 0x3f, 0x58), line_w=Pt(1))
    add_rect(sl, start_x, y, Inches(0.07), sh, fill=GOLD)
    add_text(sl, num, start_x + Inches(0.2), y + Inches(0.15),
             Inches(0.55), Inches(0.7),
             size=22, bold=True, color=RGBColor(0x2a, 0x3f, 0x58))
    add_text(sl, title, start_x + Inches(0.85), y + Inches(0.1),
             Inches(4.0), Inches(0.45),
             size=14, bold=True, color=WHITE)
    add_text(sl, body, start_x + Inches(0.85), y + Inches(0.55),
             Inches(10.0), Inches(0.48),
             size=11, color=RGBColor(0xcc, 0xd6, 0xf1))

add_text(sl, "05 / 09", W - Inches(1.5), H - Inches(0.45), Inches(1.4), Inches(0.35),
         size=9, color=RGBColor(0x44, 0x50, 0x60), align=PP_ALIGN.RIGHT)


# ============================================================
# Slide 06 – Flow
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=GRAY_BG)
add_rect(sl, 0, 0, W, Inches(0.08), fill=GOLD)

add_label(sl, "PROCESS", Inches(0.7), Inches(0.6))
add_heading(sl, "ご相談の流れ", Inches(0.7), Inches(0.95), dark_bg=False)
add_divider(sl, Inches(0.7), Inches(1.75))

steps = [
    ("1", "無料相談", "現状の課題や\n作りたいものをヒアリング"),
    ("2", "整理",     "業務・発信・導線を可視化し、\n優先順位を決める"),
    ("3", "設計",     "AIで何を効率化・制作するか\n具体化する"),
    ("4", "実装",     "ホームページ・業務フロー・\nSNS導線などを形にする"),
]

# connector line
add_rect(sl, Inches(1.9), Inches(3.4), Inches(9.5), Pt(3), fill=GOLD)

bw, bh = Inches(2.7), Inches(3.2)
gap_x = Inches(0.3)
start_x = Inches(0.6)
circle_y = Inches(2.4)
circle_r = Inches(0.7)

for i, (num, title, body) in enumerate(steps):
    cx = start_x + i * (bw + gap_x) + bw / 2

    # circle background
    add_rect(sl, cx - circle_r / 2, circle_y - circle_r / 2,
             circle_r, circle_r, fill=NAVY,
             line=GOLD, line_w=Pt(2.5))
    add_text(sl, num,
             cx - circle_r / 2, circle_y - circle_r / 2,
             circle_r, circle_r,
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    bx = start_x + i * (bw + gap_x)
    by = Inches(3.7)
    add_rect(sl, bx, by, bw, Inches(2.5), fill=WHITE,
             line=RGBColor(0xe0, 0xe4, 0xea), line_w=Pt(1))
    add_text(sl, f"STEP {num}", bx + Inches(0.15), by + Inches(0.12),
             bw - Inches(0.25), Inches(0.38),
             size=9, bold=True, color=GOLD)
    add_text(sl, title, bx + Inches(0.15), by + Inches(0.52),
             bw - Inches(0.25), Inches(0.52),
             size=16, bold=True, color=NAVY)
    add_text(sl, body, bx + Inches(0.15), by + Inches(1.1),
             bw - Inches(0.25), Inches(1.25),
             size=11, color=TEXT_MID, wrap=True)

add_text(sl, "06 / 09", W - Inches(1.5), H - Inches(0.45), Inches(1.4), Inches(0.35),
         size=9, color=TEXT_MID, align=PP_ALIGN.RIGHT)


# ============================================================
# Slide 07 – Achievements
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=WHITE)
add_rect(sl, 0, 0, W, Inches(0.08), fill=GOLD)

add_label(sl, "TRACK RECORD", Inches(0.7), Inches(0.6))
add_heading(sl, "実績・登壇歴", Inches(0.7), Inches(0.95), dark_bg=False)
add_divider(sl, Inches(0.7), Inches(1.75))

achievements = [
    ("2026年4月24日",
     "NAGOYA CONNÉCT — モデレーター登壇",
     "名古屋の女性コミュニティイベントにてモデレーターを務め、地方における男女収入格差と\n"
     "AI活用による選択肢の広がりをテーマに発信。"),
    ("2026年2月1日",
     "学び舎mom株式会社主催「集中伴走ワークショップ」— メイン講師",
     "女性起業家・経営者向けのDX・AI活用講座にて、メイン講師として登壇。\n"
     "実務で使えるAI活用の手順と型を指導。"),
    ("2025年2月・3月",
     "一般社団法人 日本小児発達子育て支援協会「VARY」— セミナー講師",
     "オンラインコミュニティ「VARY」にて、AIを活用した教材・アプリづくりをレクチャー。\n"
     "2回にわたり登壇。"),
]

aw = Inches(11.5)
ah = Inches(1.45)
start_x, start_y = Inches(0.7), Inches(2.05)
gap_y = Inches(0.22)

for i, (date, title, body) in enumerate(achievements):
    y = start_y + i * (ah + gap_y)
    add_rect(sl, start_x, y, aw, ah, fill=GRAY_BG,
             line=RGBColor(0xe0, 0xe4, 0xea), line_w=Pt(1))
    add_rect(sl, start_x, y, Inches(0.07), ah, fill=GOLD)
    add_text(sl, date, start_x + Inches(0.2), y + Inches(0.1),
             Inches(3.5), Inches(0.35),
             size=9, color=TEXT_MID)
    add_text(sl, title, start_x + Inches(0.2), y + Inches(0.42),
             aw - Inches(0.3), Inches(0.45),
             size=13, bold=True, color=NAVY)
    add_text(sl, body, start_x + Inches(0.2), y + Inches(0.87),
             aw - Inches(0.3), Inches(0.5),
             size=10, color=TEXT_MID, wrap=True)

add_text(sl, "07 / 09", W - Inches(1.5), H - Inches(0.45), Inches(1.4), Inches(0.35),
         size=9, color=TEXT_MID, align=PP_ALIGN.RIGHT)


# ============================================================
# Slide 08 – Profile
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=NAVY)
add_rect(sl, 0, 0, W, Inches(0.08), fill=GOLD)

add_label(sl, "PROFILE", Inches(0.7), Inches(0.6))

# Left panel
add_rect(sl, Inches(0.7), Inches(1.2), Inches(3.2), Inches(5.8),
         fill=NAVY_LIGHT, line=RGBColor(0x2a, 0x3f, 0x58), line_w=Pt(1))

# Profile image placeholder (gold circle)
add_rect(sl, Inches(1.35), Inches(1.5), Inches(1.9), Inches(1.9),
         fill=RGBColor(0x1e, 0x35, 0x52), line=GOLD, line_w=Pt(2))
add_text(sl, "勝目 麻希", Inches(0.7), Inches(3.5), Inches(3.2), Inches(0.55),
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "AI活用インストラクター / プロライター",
         Inches(0.7), Inches(4.05), Inches(3.2), Inches(0.45),
         size=9.5, color=GOLD, align=PP_ALIGN.CENTER)

tags = ["元銀行員", "プロライター", "名古屋拠点", "オンライン対応可"]
for j, tag in enumerate(tags):
    ty = Inches(4.65) + j * Inches(0.5)
    add_rect(sl, Inches(0.95), ty, Inches(2.7), Inches(0.38),
             fill=NAVY, line=RGBColor(0x2a, 0x3f, 0x58), line_w=Pt(1))
    add_text(sl, tag, Inches(0.95), ty + Inches(0.02), Inches(2.7), Inches(0.35),
             size=10, color=GOLD, align=PP_ALIGN.CENTER)

# Right panel
add_text(sl, "地元「名古屋」に密着したサポート",
         Inches(4.3), Inches(1.2), Inches(8.3), Inches(0.62),
         size=20, bold=True, color=WHITE)
add_divider(sl, Inches(4.3), Inches(1.9))

body1 = ("愛知県名古屋市を拠点に、オンラインはもちろん対面でのきめ細かなヒアリングにも対応。\n"
         "AIの機能説明だけでなく、貴社の課題を整理し、実際の業務・発信・集客導線にどう使うかまで一緒に考えます。")
add_text(sl, body1, Inches(4.3), Inches(2.05), Inches(8.3), Inches(1.4),
         size=11.5, color=RGBColor(0xcc, 0xd6, 0xf1), wrap=True)

body2 = ("元銀行員としての実務経験と、金融・ビジネス領域のライターとして培った言語化スキルを活かし、\n"
         "安全かつ実践的なAI活用を支援。「難しいことをわかりやすく、仕事で使える形にする」ことを大切にしています。")
add_text(sl, body2, Inches(4.3), Inches(3.6), Inches(8.3), Inches(1.4),
         size=11.5, color=RGBColor(0xcc, 0xd6, 0xf1), wrap=True)

add_text(sl, "08 / 09", W - Inches(1.5), H - Inches(0.45), Inches(1.4), Inches(0.35),
         size=9, color=RGBColor(0x44, 0x50, 0x60), align=PP_ALIGN.RIGHT)


# ============================================================
# Slide 09 – CTA
# ============================================================
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=NAVY)
add_rect(sl, 0, 0, W, Inches(0.08), fill=GOLD)
add_rect(sl, 0, H - Inches(0.08), W, Inches(0.08), fill=GOLD)

add_label(sl, "CONTACT", Inches(0.7), Inches(0.55))

add_text(sl, "まずは無料相談から。\n一緒に整えましょう。",
         Inches(0.7), Inches(0.95), Inches(12), Inches(1.6),
         size=32, bold=True, color=WHITE)
# Highlight "無料相談" — second text box in gold on top
add_text(sl, "まずは", Inches(0.7), Inches(0.95), Inches(2.2), Inches(0.75),
         size=32, bold=True, color=WHITE)
add_text(sl, "無料相談", Inches(2.85), Inches(0.95), Inches(3.5), Inches(0.75),
         size=32, bold=True, color=GOLD)
add_text(sl, "から。", Inches(6.2), Inches(0.95), Inches(2.5), Inches(0.75),
         size=32, bold=True, color=WHITE)

add_text(sl,
         "ホームページ制作、業務フロー整理、SNS発信、AI研修など、\n今の課題に合わせてご提案します。",
         Inches(0.7), Inches(1.95), Inches(12), Inches(0.9),
         size=12.5, color=RGBColor(0xcc, 0xd6, 0xf1))

# CTA box
add_rect(sl, Inches(0.7), Inches(3.0), Inches(6.5), Inches(3.5),
         fill=NAVY_LIGHT, line=RGBColor(0x2a, 0x3f, 0x58), line_w=Pt(1))
add_rect(sl, Inches(0.7), Inches(3.0), Inches(6.5), Inches(0.06), fill=GOLD)

add_text(sl, "無料相談でできること",
         Inches(0.95), Inches(3.1), Inches(6.0), Inches(0.45),
         size=11, bold=True, color=GOLD)

items = [
    "現在の課題・やりたいことのヒアリング",
    "AIで効率化できそうな業務の棚卸し",
    "最初の一歩となる具体的なアクション提案",
    "サービス内容・費用感のご案内",
]
for k, item in enumerate(items):
    iy = Inches(3.65) + k * Inches(0.56)
    add_text(sl, f"✓  {item}",
             Inches(0.95), iy, Inches(6.0), Inches(0.48),
             size=11.5, color=WHITE)

# Right side info
add_text(sl, "お問い合わせ", Inches(8.1), Inches(3.05), Inches(4.5), Inches(0.5),
         size=13, bold=True, color=GOLD)
add_text(sl, "https://makikatsume-ai.com/",
         Inches(8.1), Inches(3.65), Inches(4.5), Inches(0.45),
         size=11, color=RGBColor(0xcc, 0xd6, 0xf1))
add_text(sl, "FinEdit\n勝目 麻希（かつめ まき）\nAI活用インストラクター / プロライター\n名古屋市拠点 ／ オンライン全国対応",
         Inches(8.1), Inches(4.3), Inches(4.5), Inches(1.8),
         size=11, color=RGBColor(0xcc, 0xd6, 0xf1), wrap=True)

add_text(sl, "09 / 09", W - Inches(1.5), H - Inches(0.45), Inches(1.4), Inches(0.35),
         size=9, color=RGBColor(0x44, 0x50, 0x60), align=PP_ALIGN.RIGHT)


# ============================================================
# Save
# ============================================================
out_path = "sales-deck.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
